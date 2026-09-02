"""
tests/test_file_processor_pptx_create.py — Phase 4 of the "unused
capabilities" follow-up: python-pptx was already a dependency and
already used in actions/file_processor.py, but only to READ existing
presentations (_process_pptx's summarize/extract_text/analyze). This
tests the write direction: action='create' builds a brand-new .pptx from
a topic — the one file_processor() action that deliberately does NOT
require an existing/uploaded file, since there's nothing to upload yet.

Testing approach, per this project's own established convention: the
Gemini outline call (_generate_pptx_outline) is ALWAYS mocked — no test
here ever makes a real network/AI call. The actual pptx FILE BUILDING
(_build_pptx_from_outline via real python-pptx) is exercised for REAL
against a temp directory, the same "safe, non-destructive, so test it
for real" posture test_computer_settings.py already uses for the
clipboard round-trip — a local temp file is safe to actually create and
delete, unlike a real Gemini call or a real Desktop-launching action.

Run with:
    .venv/Scripts/python.exe -m tests.test_file_processor_pptx_create
"""
import shutil
import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock

import actions.file_processor as fp

_TMP = Path(tempfile.mkdtemp(prefix="jarvis_pptx_test_"))


def _cleanup():
    shutil.rmtree(_TMP, ignore_errors=True)


# ── _default_output_path ────────────────────────────────────────────

def test_default_output_path_is_slugged_and_timestamped_pptx() -> None:
    path = fp._default_output_path("Q3 Sales Review!!", ".pptx")
    assert path.suffix == ".pptx"
    assert "q3_sales_review" in path.stem.lower()
    print("test_default_output_path_is_slugged_and_timestamped_pptx: PASS")

def test_default_output_path_falls_back_to_home_when_no_desktop_folder() -> None:
    with patch.object(fp.Path, "is_dir", return_value=False):
        path = fp._default_output_path("topic", ".pptx")
    assert path.parent == Path.home()
    print("test_default_output_path_falls_back_to_home_when_no_desktop_folder: PASS")


# ── _generate_pptx_outline: mocked Gemini call, real JSON parsing ──────

def test_generate_pptx_outline_strips_markdown_fences_and_parses_json() -> None:
    fake_model = MagicMock()
    fake_model.generate_content.return_value = MagicMock(
        text='```json\n[{"title": "Intro", "bullets": []}, {"title": "Point 1", "bullets": ["a", "b"]}]\n```'
    )
    with patch.object(fp, "_gemini_client", return_value=fake_model):
        outline = fp._generate_pptx_outline("some topic")
    assert outline == [{"title": "Intro", "bullets": []}, {"title": "Point 1", "bullets": ["a", "b"]}]
    print("test_generate_pptx_outline_strips_markdown_fences_and_parses_json: PASS")


# ── _build_pptx_from_outline: REAL python-pptx, real temp file ─────────

def test_build_pptx_from_outline_creates_a_real_file_with_correct_slides() -> None:
    out = _TMP / "real_build_test.pptx"
    outline = [
        {"title": "Welcome", "bullets": []},
        {"title": "Key Points", "bullets": ["First point", "Second point", "Third point"]},
    ]
    fp._build_pptx_from_outline(out, outline)
    assert out.exists()

    from pptx import Presentation
    prs = Presentation(str(out))
    assert len(list(prs.slides)) == 2
    titles = [s.shapes.title.text for s in prs.slides]
    assert titles == ["Welcome", "Key Points"]
    body_text = prs.slides[1].placeholders[1].text_frame.text
    assert "First point" in body_text and "Second point" in body_text and "Third point" in body_text
    print("test_build_pptx_from_outline_creates_a_real_file_with_correct_slides: PASS")


# ── _create_pptx: full flow, mocked outline generation only ────────────

def test_create_pptx_with_explicit_slides_skips_the_gemini_call_entirely() -> None:
    out = _TMP / "explicit_slides.pptx"
    slides = [{"title": "A", "bullets": ["x"]}]
    with patch.object(fp, "_generate_pptx_outline") as m_gen:
        result = fp._create_pptx(out, "", {"slides": slides})
    m_gen.assert_not_called()
    assert out.exists()
    assert "1 slide" in result
    print("test_create_pptx_with_explicit_slides_skips_the_gemini_call_entirely: PASS")

def test_create_pptx_with_only_a_topic_uses_the_generated_outline() -> None:
    out = _TMP / "from_topic.pptx"
    fake_outline = [{"title": "Volcanoes", "bullets": []}, {"title": "Types", "bullets": ["Shield", "Composite"]}]
    with patch.object(fp, "_generate_pptx_outline", return_value=fake_outline) as m_gen:
        result = fp._create_pptx(out, "volcanoes", {})
    m_gen.assert_called_once_with("volcanoes")
    assert out.exists()
    assert "2 slide" in result
    print("test_create_pptx_with_only_a_topic_uses_the_generated_outline: PASS")

def test_create_pptx_falls_back_to_the_filename_as_topic_when_no_instruction_given() -> None:
    # No instruction/topic passed — _create_pptx falls back to the
    # target filename itself as the topic rather than refusing outright,
    # since a path was still explicitly chosen.
    out = _TMP / "quarterly_budget_review.pptx"
    with patch.object(fp, "_generate_pptx_outline", return_value=[{"title": "X", "bullets": []}]) as m_gen:
        fp._create_pptx(out, "", {})
    m_gen.assert_called_once_with("quarterly budget review")
    print("test_create_pptx_falls_back_to_the_filename_as_topic_when_no_instruction_given: PASS")

def test_create_pptx_with_no_topic_and_no_filename_to_fall_back_to_calls_nothing() -> None:
    with patch.object(fp, "_generate_pptx_outline") as m_gen:
        result = fp._create_pptx(Path(""), "", {})
    m_gen.assert_not_called()
    assert "No topic" in result
    print("test_create_pptx_with_no_topic_and_no_filename_to_fall_back_to_calls_nothing: PASS")

def test_create_pptx_never_crashes_when_outline_generation_fails() -> None:
    out = _TMP / "should_not_exist.pptx"
    with patch.object(fp, "_generate_pptx_outline", side_effect=RuntimeError("bad json")):
        result = fp._create_pptx(out, "some topic", {})
    assert "Could not generate" in result
    assert not out.exists()
    print("test_create_pptx_never_crashes_when_outline_generation_fails: PASS")

def test_create_pptx_never_crashes_when_the_pptx_build_itself_fails() -> None:
    out = _TMP / "build_failure.pptx"
    with patch.object(fp, "_generate_pptx_outline", return_value=[{"title": "X", "bullets": []}]), \
         patch.object(fp, "_build_pptx_from_outline", side_effect=OSError("disk full")):
        result = fp._create_pptx(out, "some topic", {})
    assert "Could not build" in result
    print("test_create_pptx_never_crashes_when_the_pptx_build_itself_fails: PASS")


# ── file_processor(): top-level dispatch wiring ─────────────────────────

def test_file_processor_create_action_uses_default_path_when_none_given() -> None:
    fake_path = _TMP / "default_path_used.pptx"
    with patch.object(fp, "_default_output_path", return_value=fake_path) as m_default, \
         patch.object(fp, "_create_pptx", return_value="[ok]") as m_create:
        result = fp.file_processor({"action": "create", "instruction": "quarterly review"})
    m_default.assert_called_once_with("quarterly review", ".pptx")
    m_create.assert_called_once_with(fake_path, "quarterly review", {"action": "create", "instruction": "quarterly review"})
    assert result == "[ok]"
    print("test_file_processor_create_action_uses_default_path_when_none_given: PASS")

def test_file_processor_create_action_respects_an_explicit_file_path() -> None:
    explicit = str(_TMP / "explicit.pptx")
    with patch.object(fp, "_default_output_path") as m_default, \
         patch.object(fp, "_create_pptx", return_value="[ok]") as m_create:
        fp.file_processor({"action": "create", "file_path": explicit, "instruction": "topic"})
    m_default.assert_not_called()
    m_create.assert_called_once()
    called_path = m_create.call_args[0][0]
    assert str(called_path) == explicit
    print("test_file_processor_create_action_respects_an_explicit_file_path: PASS")

def test_file_processor_create_action_rejects_non_pptx_formats_honestly() -> None:
    with patch.object(fp, "_create_pptx") as m_create:
        result = fp.file_processor({"action": "create", "format": "docx", "instruction": "topic"})
    m_create.assert_not_called()
    assert "only supports pptx" in result
    print("test_file_processor_create_action_rejects_non_pptx_formats_honestly: PASS")

def test_file_processor_non_create_action_on_a_missing_file_is_unchanged() -> None:
    # Regression guard: the existing "process an uploaded file" behavior
    # for a genuinely missing file must be untouched by the new branch.
    result = fp.file_processor({"action": "summarize", "file_path": str(_TMP / "does_not_exist.pptx")})
    assert result.startswith("File not found")
    print("test_file_processor_non_create_action_on_a_missing_file_is_unchanged: PASS")


if __name__ == "__main__":
    try:
        test_default_output_path_is_slugged_and_timestamped_pptx()
        test_default_output_path_falls_back_to_home_when_no_desktop_folder()
        test_generate_pptx_outline_strips_markdown_fences_and_parses_json()
        test_build_pptx_from_outline_creates_a_real_file_with_correct_slides()
        test_create_pptx_with_explicit_slides_skips_the_gemini_call_entirely()
        test_create_pptx_with_only_a_topic_uses_the_generated_outline()
        test_create_pptx_falls_back_to_the_filename_as_topic_when_no_instruction_given()
        test_create_pptx_with_no_topic_and_no_filename_to_fall_back_to_calls_nothing()
        test_create_pptx_never_crashes_when_outline_generation_fails()
        test_create_pptx_never_crashes_when_the_pptx_build_itself_fails()
        test_file_processor_create_action_uses_default_path_when_none_given()
        test_file_processor_create_action_respects_an_explicit_file_path()
        test_file_processor_create_action_rejects_non_pptx_formats_honestly()
        test_file_processor_non_create_action_on_a_missing_file_is_unchanged()
        print("\nAll file_processor pptx-create tests passed.")
    finally:
        _cleanup()
